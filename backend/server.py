from contextlib import contextmanager
from datetime import datetime, timedelta, timezone
import bcrypt
from bs4 import BeautifulSoup
import certifi
from docx import Document as DocxDocument
from dotenv import load_dotenv
from fastapi import APIRouter, Depends, FastAPI, File, Header, HTTPException, Request, UploadFile
from fastapi.responses import FileResponse
import hashlib
import hmac
import io
import json
import logging
import os
from pathlib import Path
from pydantic import BaseModel, EmailStr
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import pg8000
from pypdf import PdfReader
import pdfplumber
from pptx import Presentation as PptxPresentation
import jwt as pyjwt
import razorpay
import re
import requests
import ssl
from starlette.middleware.cors import CORSMiddleware
import tempfile
from typing import List, Optional
from urllib.parse import unquote, urlparse
import uuid

from emergentintegrations.llm.chat import FileContentWithMimeType, LlmChat, UserMessage

BASE_DIR = Path(__file__).resolve().parent.parent
BACKEND_DIR = BASE_DIR / "backend"
FRONTEND_BUILD_DIR = BASE_DIR / "frontend" / "build"
load_dotenv(BACKEND_DIR / ".env")

DATABASE_URL = os.environ.get("SUPABASE_DB_URL") or os.environ.get("DATABASE_URL")
EMERGENT_LLM_KEY = os.environ["EMERGENT_LLM_KEY"]
RAZORPAY_KEY_ID = os.environ["RAZORPAY_KEY_ID"]
RAZORPAY_KEY_SECRET = os.environ["RAZORPAY_KEY_SECRET"]
JWT_SECRET = os.environ["JWT_SECRET"]

ROAST_PRICE_PAISE = 1000  # Rs 10
PREMIUM_UPGRADE_PAISE = 2900  # Rs 29

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(name)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

if not DATABASE_URL:
    raise RuntimeError("SUPABASE_DB_URL or DATABASE_URL must be set")

SCHEMA_SQL = """
CREATE TABLE IF NOT EXISTS users (
    id TEXT PRIMARY KEY,
    email TEXT NOT NULL UNIQUE,
    name TEXT NOT NULL,
    password_hash TEXT NOT NULL,
    used_free_roast BOOLEAN NOT NULL DEFAULT FALSE,
    paid_roasts_balance INTEGER NOT NULL DEFAULT 0,
    created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
);

CREATE TABLE IF NOT EXISTS roasts (
    id TEXT PRIMARY KEY,
    user_id TEXT NOT NULL REFERENCES users(id) ON DELETE CASCADE,
    user_name TEXT NOT NULL,
    startup_name TEXT NOT NULL DEFAULT '',
    idea TEXT NOT NULL,
    score INTEGER NOT NULL,
    one_liner TEXT NOT NULL,
    callouts JSONB NOT NULL,
    fixes JSONB NOT NULL,
    verdict_title TEXT NOT NULL,
    is_premium BOOLEAN NOT NULL DEFAULT FALSE,
    competitors JSONB,
    tam_analysis TEXT,
    india_rating INTEGER,
    global_rating INTEGER,
    created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
);

CREATE INDEX IF NOT EXISTS roasts_user_created_idx ON roasts (user_id, created_at DESC);
CREATE INDEX IF NOT EXISTS roasts_leaderboard_idx ON roasts (score ASC, created_at DESC);

CREATE TABLE IF NOT EXISTS payments (
    order_id TEXT PRIMARY KEY,
    user_id TEXT NOT NULL REFERENCES users(id) ON DELETE CASCADE,
    amount INTEGER NOT NULL,
    status TEXT NOT NULL,
    payment_id TEXT,
    payment_type TEXT NOT NULL DEFAULT 'roast_credit',
    roast_id TEXT REFERENCES roasts(id) ON DELETE SET NULL,
    created_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
    paid_at TIMESTAMPTZ
);

CREATE TABLE IF NOT EXISTS feedback (
    id TEXT PRIMARY KEY,
    user_id TEXT REFERENCES users(id) ON DELETE SET NULL,
    type TEXT NOT NULL,
    content TEXT NOT NULL,
    email TEXT,
    created_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
);
"""


@contextmanager
def get_db():
    parsed = urlparse(DATABASE_URL)
    try:
        conn = pg8000.dbapi.connect(
            user=unquote(parsed.username or ""),
            password=unquote(parsed.password or ""),
            host=parsed.hostname,
            port=parsed.port or 5432,
            database=(parsed.path or "/postgres").lstrip("/"),
            ssl_context=ssl._create_unverified_context(),
        )
    except Exception:
        logger.exception(
            "Failed to connect to Postgres host=%s port=%s database=%s",
            parsed.hostname,
            parsed.port or 5432,
            (parsed.path or "/postgres").lstrip("/"),
        )
        raise
    try:
        yield _ConnectionWrapper(conn)
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        conn.close()


class _CursorWrapper:
    def __init__(self, cursor):
        self._cursor = cursor

    def execute(self, *args, **kwargs):
        return self._cursor.execute(*args, **kwargs)

    def fetchone(self):
        row = self._cursor.fetchone()
        return self._as_dict(row)

    def fetchall(self):
        rows = self._cursor.fetchall()
        return [self._as_dict(row) for row in rows]

    def _as_dict(self, row):
        if row is None:
            return None
        return {desc[0]: value for desc, value in zip(self._cursor.description, row)}

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        self._cursor.close()
        return False


class _ConnectionWrapper:
    def __init__(self, conn):
        self._conn = conn

    def cursor(self):
        return _CursorWrapper(self._conn.cursor())


def init_db():
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute(SCHEMA_SQL)


def _row_to_user_public(row: dict) -> dict:
    return {
        "id": row["id"],
        "email": row["email"],
        "name": row["name"],
        "used_free_roast": bool(row.get("used_free_roast", False)),
        "paid_roasts_balance": int(row.get("paid_roasts_balance", 0)),
    }


def _row_to_roast(row: dict) -> dict:
    is_premium = bool(row.get("is_premium", False))
    return {
        "id": row["id"],
        "user_id": row["user_id"],
        "user_name": row["user_name"],
        "startup_name": row["startup_name"],
        "idea": row["idea"],
        "score": int(row["score"]),
        "one_liner": row["one_liner"],
        "callouts": list(row["callouts"] or []),
        "fixes": list(row["fixes"] or []),
        "verdict_title": row["verdict_title"],
        "created_at": row["created_at"].astimezone(timezone.utc).isoformat(),
        "is_premium": is_premium,
        "competitors": list(row["competitors"] or []) if is_premium else None,
        "tam_analysis": row["tam_analysis"] if is_premium else None,
        "india_rating": int(row["india_rating"]) if is_premium and row.get("india_rating") is not None else None,
        "global_rating": int(row["global_rating"]) if is_premium and row.get("global_rating") is not None else None,
    }


async def _backfill_premium_data(roast_id: str, startup_name: str, idea: str):
    try:
        roast_data = await generate_roast_ai(startup_name, idea)
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    """
                    UPDATE roasts
                    SET competitors = %s::jsonb,
                        tam_analysis = %s,
                        india_rating = %s,
                        global_rating = %s
                    WHERE id = %s
                    """,
                    (
                        json.dumps(roast_data["competitors"]),
                        roast_data["tam_analysis"],
                        roast_data["india_rating"],
                        roast_data["global_rating"],
                        roast_id,
                    ),
                )
        return roast_data
    except Exception as e:
        logger.error(f"Failed to backfill premium data for {roast_id}: {e}")
        return None


init_db()

rzp_client = razorpay.Client(auth=(RAZORPAY_KEY_ID, RAZORPAY_KEY_SECRET))

app = FastAPI()
api_router = APIRouter(prefix="/api")


class RegisterInput(BaseModel):
    email: EmailStr
    password: str
    name: str


class LoginInput(BaseModel):
    email: EmailStr
    password: str


class UserPublic(BaseModel):
    id: str
    email: str
    name: str
    used_free_roast: bool
    paid_roasts_balance: int


class RoastOut(BaseModel):
    id: str
    user_id: str
    user_name: str
    startup_name: str
    idea: str
    score: int
    one_liner: str
    callouts: List[str]
    fixes: List[str]
    verdict_title: str
    created_at: str
    is_premium: bool = False
    competitors: Optional[List[str]] = None
    tam_analysis: Optional[str] = None
    india_rating: Optional[int] = None
    global_rating: Optional[int] = None


class CreateOrderOut(BaseModel):
    order_id: str
    amount: int
    currency: str
    key_id: str


class VerifyPaymentInput(BaseModel):
    razorpay_order_id: str
    razorpay_payment_id: str
    razorpay_signature: str


class FeedbackInput(BaseModel):
    type: str  # 'feedback' or 'feature_request'
    content: str
    email: Optional[str] = None


async def get_optional_user(authorization: Optional[str] = Header(None)):
    if not authorization or not authorization.lower().startswith("bearer "):
        return None
    token = authorization.split(" ", 1)[1].strip()
    try:
        payload = pyjwt.decode(token, JWT_SECRET, algorithms=["HS256"])
        user_id = payload.get("sub")
        return _fetch_user_by_id(user_id)
    except Exception:
        return None


def hash_password(pw: str) -> str:
    return bcrypt.hashpw(pw.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")


def verify_password(pw: str, hashed: str) -> bool:
    return bcrypt.checkpw(pw.encode("utf-8"), hashed.encode("utf-8"))


def create_jwt(user_id: str) -> str:
    payload = {
        "sub": user_id,
        "iat": datetime.now(timezone.utc),
        "exp": datetime.now(timezone.utc) + timedelta(days=30),
    }
    return pyjwt.encode(payload, JWT_SECRET, algorithm="HS256")


def _fetch_user_by_id(user_id: str) -> Optional[dict]:
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT id, email, name, used_free_roast, paid_roasts_balance
                FROM users
                WHERE id = %s
                """,
                (user_id,),
            )
            return cur.fetchone()


async def get_current_user(authorization: Optional[str] = Header(None)):
    if not authorization or not authorization.lower().startswith("bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid Authorization header")
    token = authorization.split(" ", 1)[1].strip()
    try:
        payload = pyjwt.decode(token, JWT_SECRET, algorithms=["HS256"])
        user_id = payload.get("sub")
    except Exception as exc:
        raise HTTPException(status_code=401, detail="Invalid or expired token") from exc
    user = _fetch_user_by_id(user_id)
    if not user:
        raise HTTPException(status_code=401, detail="User not found")
    return user


def _extract_pdf_text(data: bytes) -> str:
    # Try pdfplumber first
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            text = ""
            for page in pdf.pages[:30]:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            if text.strip():
                return text.strip()
    except Exception as e:
        logger.warning(f"pdfplumber failed to extract text: {e}")

    # Fallback to pypdf
    try:
        reader = PdfReader(io.BytesIO(data))
        chunks = []
        for page in reader.pages[:30]:
            try:
                chunks.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n".join(chunks).strip()
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not read PDF: {exc}") from exc


def _extract_docx_text(data: bytes) -> str:
    try:
        doc = DocxDocument(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs).strip()
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not read DOCX: {exc}") from exc


def _extract_pptx_text(data: bytes) -> str:
    try:
        prs = PptxPresentation(io.BytesIO(data))
        chunks = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                chunks.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(chunks).strip()
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not read PPTX: {exc}") from exc


TEXTY_EXTS = {".txt", ".md", ".csv", ".tsv", ".rtf", ".html", ".htm", ".json", ".log", ".xml", ".yml", ".yaml"}
IMAGE_MIMES = {"image/png", "image/jpeg", "image/jpg", "image/webp", "image/gif", "image/heic", "image/heif"}


def _extract_file_content(filename: str, content_type: str, data: bytes):
    name = (filename or "").lower()
    ext = "." + name.rsplit(".", 1)[-1] if "." in name else ""
    ctype = (content_type or "").lower()

    if ext == ".pdf" or ctype == "application/pdf":
        return (_extract_pdf_text(data), None, None)

    if ext == ".docx" or ctype in {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}:
        return (_extract_docx_text(data), None, None)

    if ext == ".pptx" or ctype in {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}:
        return (_extract_pptx_text(data), None, None)

    if ctype in IMAGE_MIMES or ext in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".heic", ".heif"}:
        suffix = ext or ".png"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(data)
        tmp.close()
        mime = ctype if ctype in IMAGE_MIMES else ("image/jpeg" if suffix in {".jpg", ".jpeg"} else "image/png")
        return (None, tmp.name, mime)

    if ext in TEXTY_EXTS or (ctype.startswith("text/") if ctype else False):
        try:
            txt = data.decode("utf-8", errors="ignore")
            if ext in {".html", ".htm"}:
                soup = BeautifulSoup(txt, "html.parser")
                for tag in soup(["script", "style"]):
                    tag.decompose()
                txt = " ".join(soup.get_text(" ").split())
            return (txt[:20000], None, None)
        except Exception:
            pass

    try:
        txt = data.decode("utf-8")
        if txt.strip():
            return (txt[:20000], None, None)
    except Exception:
        pass

    raise HTTPException(
        status_code=400,
        detail=f"Unsupported file type: {ext or ctype or 'unknown'}. Try PDF, PPTX, DOCX, images, or text files.",
    )


def _extract_url_text(url: str) -> str:
    if not re.match(r"^https?://", url, flags=re.I):
        url = "https://" + url
    try:
        resp = requests.get(
            url,
            timeout=15,
            headers={"User-Agent": "Mozilla/5.0 (Roastmaster/1.0)"},
            allow_redirects=True,
        )
        resp.raise_for_status()
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not fetch URL: {exc}") from exc
    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "noscript", "svg"]):
        tag.decompose()
    title = soup.title.string.strip() if soup.title and soup.title.string else ""
    text = " ".join(soup.get_text(" ").split())
    combined = f"Page Title: {title}\n\n{text}" if title else text
    return combined[:8000]


ROAST_SYSTEM = """You are THE ROASTMASTER — an AI angel investor who has seen 10,000 pitch decks and has zero patience for fluff. You review startup ideas like a brutally honest investor who actually cares enough to tell the truth. You are savage, witty, and surgically precise — but never hateful. You find the REAL, SPECIFIC flaws — market reality, execution, TAM, moat, differentiation, founder delusion.

Output ONLY valid JSON in this exact schema, nothing else, no prefix, no suffix, no markdown fences:
{
  "score": <integer 1-10, where 1 = absolute trash and 10 = unicorn>,
  "verdict_title": "<a single brutal headline verdict, max 8 words, all caps>",
  "one_liner": "<one biting sentence summing up why this idea is cooked (max 25 words)>",
  "callouts": [
    "<brutal specific callout 1 — reference actual details from their idea>",
    "<callout 2>",
    "<callout 3>",
    "<callout 4>",
    "<callout 5>"
  ],
  "fixes": [
    "<sharp actionable fix 1>",
    "<fix 2>",
    "<fix 3>",
    "<fix 4>",
    "<fix 5>"
  ],
  "competitors": [
    "<name of real competitor 1>",
    "<name of real competitor 2>",
    "<name of real competitor 3>"
  ],
  "tam_analysis": "<1-2 sentences brutally explaining why their market size is either a lie or impossible to capture>",
  "india_rating": <integer 0-100, probability of survival in the Indian market>,
  "global_rating": <integer 0-100, probability of global scaling success>
}

Rules:
- Callouts must be SPECIFIC to their idea, not generic advice.
- Use dry wit, irony, and dark startup humor. Think: Ricky Gervais roasting a TechCrunch pitch.
- Do NOT be supportive. Do NOT say 'great idea but'. Do NOT hedge.
- Each callout should be 1-2 sentences, punchy.
- Fixes should be specific, not 'do more research'.
- Score strictly: most ideas deserve 2-5. Unicorns are rare.
- Competitors MUST be real existing companies or startups.
- Ratings should be realistic. If it's a hyper-local Indian idea, global_rating should be low.
- Do not use emoji.
- Return only JSON. No commentary."""


async def generate_roast_ai(startup_name: str, idea: str, file_attachments: Optional[List[dict]] = None) -> dict:
    session_id = str(uuid.uuid4())
    chat = (
        LlmChat(api_key=EMERGENT_LLM_KEY, session_id=session_id, system_message=ROAST_SYSTEM)
        .with_model("gemini", "gemini-3-flash-preview")
    )
    user_text = f"Startup Name: {startup_name or 'Unnamed'}\n\nIdea / Materials:\n{idea}\n\nRoast it. Return ONLY the JSON."

    file_contents = None
    if file_attachments:
        file_contents = [FileContentWithMimeType(mime_type=item["mime"], file_path=item["path"]) for item in file_attachments]

    msg = UserMessage(text=user_text, file_contents=file_contents) if file_contents else UserMessage(text=user_text)
    response = await chat.send_message(msg)
    text = (response or "").strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z]*", "", text).strip()
        text = text.rstrip("`").strip()
        if text.endswith("```"):
            text = text[:-3].strip()
    try:
        data = json.loads(text)
    except Exception:
        match = re.search(r"\{[\s\S]*\}", text)
        if not match:
            raise HTTPException(status_code=502, detail="AI returned invalid response")
        data = json.loads(match.group(0))

    score = max(1, min(10, int(data.get("score", 3))))
    callouts = [str(x).strip() for x in (data.get("callouts") or [])][:5]
    fixes = [str(x).strip() for x in (data.get("fixes") or [])][:5]
    while len(callouts) < 5:
        callouts.append("The silence speaks louder than your pitch.")
    while len(fixes) < 5:
        fixes.append("Start over. Talk to 50 real users first.")

    def safe_int(v, default=0):
        try:
            return int(v)
        except (ValueError, TypeError):
            return default

    return {
        "score": score,
        "verdict_title": str(data.get("verdict_title", "DELUSION DETECTED"))[:120].upper(),
        "one_liner": str(data.get("one_liner", "This idea needs a funeral, not a Series A."))[:300],
        "callouts": callouts,
        "fixes": fixes,
        "competitors": [str(c).strip() for c in (data.get("competitors") or [])][:3],
        "tam_analysis": str(data.get("tam_analysis", "")).strip(),
        "india_rating": max(0, min(100, safe_int(data.get("india_rating"), 0))),
        "global_rating": max(0, min(100, safe_int(data.get("global_rating"), 0))),
    }


@api_router.post("/auth/register")
async def register(payload: RegisterInput):
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id FROM users WHERE email = %s", (payload.email.lower(),))
            if cur.fetchone():
                raise HTTPException(status_code=400, detail="Email already registered")

            user_id = str(uuid.uuid4())
            cur.execute(
                """
                INSERT INTO users (id, email, name, password_hash, used_free_roast, paid_roasts_balance, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                RETURNING id, email, name, used_free_roast, paid_roasts_balance
                """,
                (
                    user_id,
                    payload.email.lower(),
                    payload.name.strip()[:60] or "Anon",
                    hash_password(payload.password),
                    False,
                    0,
                    datetime.now(timezone.utc),
                ),
            )
            user = cur.fetchone()
    token = create_jwt(user["id"])
    return {"token": token, "user": _row_to_user_public(user)}


@api_router.post("/auth/login")
async def login(payload: LoginInput):
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT * FROM users WHERE email = %s", (payload.email.lower(),))
            user = cur.fetchone()
    if not user or not verify_password(payload.password, user["password_hash"]):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    token = create_jwt(user["id"])
    return {"token": token, "user": _row_to_user_public(user)}


@api_router.get("/auth/me")
async def me(current=Depends(get_current_user)):
    return _row_to_user_public(current)


@api_router.post("/roast/generate", response_model=RoastOut)
async def generate_roast(request: Request, current=Depends(get_current_user)):
    content_type = request.headers.get("content-type", "")
    idea = ""
    startup_name = ""
    source_url = ""
    uploaded_file = None

    if "application/json" in content_type:
        payload = await request.json()
        idea = str(payload.get("idea", ""))
        startup_name = str(payload.get("startup_name", ""))
        source_url = str(payload.get("source_url", ""))
    else:
        form = await request.form()
        idea = str(form.get("idea", ""))
        startup_name = str(form.get("startup_name", ""))
        source_url = str(form.get("source_url", ""))
        maybe_file = form.get("pdf_file")
        if isinstance(maybe_file, UploadFile):
            uploaded_file = maybe_file

    if len(idea.strip()) > 4000:
        raise HTTPException(status_code=400, detail="Idea too long (max 4000 chars).")

    parts = []
    file_attachments = []
    temp_paths_to_clean = []

    if idea and idea.strip():
        parts.append(idea.strip())
    if source_url and source_url.strip():
        parts.append("\n---\nFrom URL:\n" + _extract_url_text(source_url.strip()))
    if uploaded_file is not None:
        data = await uploaded_file.read()
        if len(data) > 15 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="File too large (max 15 MB)")
        extracted_text, tmp_path, tmp_mime = _extract_file_content(uploaded_file.filename or "upload", uploaded_file.content_type or "", data)
        if extracted_text:
            parts.append(f"\n---\nFrom uploaded file ({uploaded_file.filename}):\n" + extracted_text[:8000])
        if tmp_path:
            file_attachments.append({"path": tmp_path, "mime": tmp_mime})
            temp_paths_to_clean.append(tmp_path)
            parts.append(f"\n---\nUploaded image attached: {uploaded_file.filename}. Analyse it as part of the pitch.")

    combined = "\n".join(parts).strip()
    if not combined and not file_attachments:
        raise HTTPException(status_code=400, detail="No content detected to roast. Please upload a valid file or add a text description.")
    if not combined:
        combined = "See attached image(s)."
    if len(combined) < 15 and not file_attachments:
        raise HTTPException(status_code=400, detail="Not enough content to roast. Please add some more text (min 15 characters) or a more detailed file.")
    if len(combined) > 12000:
        combined = combined[:12000]

    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT used_free_roast, paid_roasts_balance FROM users WHERE id = %s", (current["id"],))
            entitlement = cur.fetchone()
    if not entitlement:
        raise HTTPException(status_code=401, detail="User not found")

    try:
        roast_data = await generate_roast_ai(startup_name or "", combined, file_attachments=file_attachments)
    finally:
        for path in temp_paths_to_clean:
            try:
                os.unlink(path)
            except Exception:
                pass

    roast_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc)
    stored_idea = idea.strip() if idea and idea.strip() else combined[:2000]

    with get_db() as conn:
        with conn.cursor() as cur:
            if not entitlement["used_free_roast"]:
                cur.execute(
                    """
                    UPDATE users
                    SET used_free_roast = TRUE
                    WHERE id = %s AND used_free_roast = FALSE
                    RETURNING id
                    """,
                    (current["id"],),
                )
            elif entitlement["paid_roasts_balance"] > 0:
                cur.execute(
                    """
                    UPDATE users
                    SET paid_roasts_balance = paid_roasts_balance - 1
                    WHERE id = %s AND paid_roasts_balance > 0
                    RETURNING id
                    """,
                    (current["id"],),
                )
            else:
                raise HTTPException(status_code=402, detail="Free roast consumed. Pay Rs 10 for another savage review.")

            if not cur.fetchone():
                raise HTTPException(status_code=402, detail="Roast credit is no longer available. Please try again.")

            cur.execute(
                """
                INSERT INTO roasts (
                    id, user_id, user_name, startup_name, idea, score, one_liner,
                    callouts, fixes, verdict_title, competitors, tam_analysis,
                    india_rating, global_rating, created_at
                )
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s::jsonb, %s::jsonb, %s, %s::jsonb, %s, %s, %s, %s)
                RETURNING *
                """,
                (
                    roast_id,
                    current["id"],
                    current["name"],
                    (startup_name or "").strip()[:80],
                    stored_idea,
                    roast_data["score"],
                    roast_data["one_liner"],
                    json.dumps(roast_data["callouts"]),
                    json.dumps(roast_data["fixes"]),
                    roast_data["verdict_title"],
                    json.dumps(roast_data.get("competitors", [])),
                    roast_data.get("tam_analysis", ""),
                    roast_data.get("india_rating", 0),
                    roast_data.get("global_rating", 0),
                    now,
                ),
            )
            roast_row = cur.fetchone()
    return _row_to_roast(roast_row)


@api_router.get("/roast/{roast_id}", response_model=RoastOut)
async def get_roast(roast_id: str):
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT * FROM roasts WHERE id = %s", (roast_id,))
            roast = cur.fetchone()
    if not roast:
        raise HTTPException(status_code=404, detail="Roast not found")
    
    # Auto-backfill if premium but missing data
    if roast.get("is_premium") and (not roast.get("competitors") or not roast.get("tam_analysis")):
        await _backfill_premium_data(roast["id"], roast["startup_name"], roast["idea"])
        # Fetch again after backfill
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute("SELECT * FROM roasts WHERE id = %s", (roast_id,))
                roast = cur.fetchone()

    return _row_to_roast(roast)


@api_router.get("/roasts/my")
async def my_roasts(current=Depends(get_current_user)):
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT * FROM roasts WHERE user_id = %s ORDER BY created_at DESC LIMIT 100", (current["id"],))
            rows = cur.fetchall()
    return [_row_to_roast(row) for row in rows]


@api_router.get("/roasts/leaderboard")
async def leaderboard():
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT id, score, one_liner, verdict_title, user_name, created_at
                FROM roasts
                ORDER BY score ASC, created_at DESC
                LIMIT 50
                """
            )
            rows = cur.fetchall()
    return [
        {
            "id": row["id"],
            "score": int(row["score"]),
            "one_liner": row["one_liner"],
            "verdict_title": row["verdict_title"],
            "user_name": row["user_name"],
            "created_at": row["created_at"].astimezone(timezone.utc).isoformat(),
        }
        for row in rows
    ]


@api_router.post("/payment/create-order", response_model=CreateOrderOut)
async def create_order(current=Depends(get_current_user)):
    short_id = current["id"][:8]
    receipt = f"roast_{short_id}_{int(datetime.now(timezone.utc).timestamp())}"[:40]
    order = rzp_client.order.create(
        {
            "amount": ROAST_PRICE_PAISE,
            "currency": "INR",
            "receipt": receipt,
            "payment_capture": 1,
            "notes": {"user_id": current["id"], "product": "startup_roast"},
        }
    )
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO payments (order_id, user_id, amount, status, payment_type, created_at)
                VALUES (%s, %s, %s, %s, %s, %s)
                ON CONFLICT (order_id) DO UPDATE
                SET user_id = EXCLUDED.user_id,
                    amount = EXCLUDED.amount,
                    status = EXCLUDED.status,
                    payment_type = EXCLUDED.payment_type
                """,
                (order["id"], current["id"], ROAST_PRICE_PAISE, "created", "roast_credit", datetime.now(timezone.utc)),
            )
    return {"order_id": order["id"], "amount": ROAST_PRICE_PAISE, "currency": "INR", "key_id": RAZORPAY_KEY_ID}


@api_router.post("/roast/{roast_id}/premium-order", response_model=CreateOrderOut)
async def create_premium_order(roast_id: str, current=Depends(get_current_user)):
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT user_id, is_premium FROM roasts WHERE id = %s", (roast_id,))
            roast = cur.fetchone()
    if not roast:
        raise HTTPException(status_code=404, detail="Roast not found")
    if roast["user_id"] != current["id"]:
        raise HTTPException(status_code=403, detail="Not your roast")
    if roast["is_premium"]:
        raise HTTPException(status_code=400, detail="Already premium")

    receipt = f"prem_{roast_id[:8]}_{int(datetime.now(timezone.utc).timestamp())}"[:40]
    order = rzp_client.order.create(
        {
            "amount": PREMIUM_UPGRADE_PAISE,
            "currency": "INR",
            "receipt": receipt,
            "payment_capture": 1,
            "notes": {"user_id": current["id"], "roast_id": roast_id, "product": "premium_upgrade"},
        }
    )
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO payments (order_id, user_id, roast_id, amount, status, payment_type, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                """,
                (order["id"], current["id"], roast_id, PREMIUM_UPGRADE_PAISE, "created", "premium_upgrade", datetime.now(timezone.utc)),
            )
    return {"order_id": order["id"], "amount": PREMIUM_UPGRADE_PAISE, "currency": "INR", "key_id": RAZORPAY_KEY_ID}


@api_router.post("/payment/verify")
async def verify_payment(payload: VerifyPaymentInput, current=Depends(get_current_user)):
    body = f"{payload.razorpay_order_id}|{payload.razorpay_payment_id}"
    expected = hmac.new(RAZORPAY_KEY_SECRET.encode("utf-8"), body.encode("utf-8"), hashlib.sha256).hexdigest()
    if not hmac.compare_digest(expected, payload.razorpay_signature):
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    "UPDATE payments SET status = %s WHERE order_id = %s",
                    ("invalid_signature", payload.razorpay_order_id),
                )
        raise HTTPException(status_code=400, detail="Payment signature mismatch")

    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT * FROM payments WHERE order_id = %s", (payload.razorpay_order_id,))
            payment = cur.fetchone()
            if not payment or payment["user_id"] != current["id"]:
                raise HTTPException(status_code=404, detail="Order not found")
            if payment.get("status") == "paid":
                cur.execute(
                    """
                    SELECT id, email, name, used_free_roast, paid_roasts_balance
                    FROM users
                    WHERE id = %s
                    """,
                    (current["id"],),
                )
                user = cur.fetchone()
                return {"status": "already_credited", "user": _row_to_user_public(user)}

            cur.execute(
                """
                UPDATE payments
                SET status = %s, payment_id = %s, paid_at = %s
                WHERE order_id = %s
                """,
                ("paid", payload.razorpay_payment_id, datetime.now(timezone.utc), payload.razorpay_order_id),
            )

            if payment.get("payment_type") == "premium_upgrade" and payment.get("roast_id"):
                roast_id = payment["roast_id"]
                cur.execute("UPDATE roasts SET is_premium = TRUE WHERE id = %s", (roast_id,))
                
                # Check if we need to backfill data immediately
                cur.execute("SELECT idea, startup_name, competitors, tam_analysis FROM roasts WHERE id = %s", (roast_id,))
                r = cur.fetchone()
                if r and (not r.get("competitors") or not r.get("tam_analysis") or r.get("tam_analysis") == ""):
                    await _backfill_premium_data(roast_id, r["startup_name"], r["idea"])
                # Fetch user again to return consistent response
                cur.execute("SELECT id, email, name, used_free_roast, paid_roasts_balance FROM users WHERE id = %s", (current["id"],))
                user = cur.fetchone()
            else:
                cur.execute(
                    """
                    UPDATE users
                    SET paid_roasts_balance = paid_roasts_balance + 1
                    WHERE id = %s
                    RETURNING id, email, name, used_free_roast, paid_roasts_balance
                    """,
                    (current["id"],),
                )
                user = cur.fetchone()
    return {"status": "success", "user": _row_to_user_public(user)}


@api_router.post("/feedback")
async def submit_feedback(payload: FeedbackInput, current=Depends(get_optional_user)):
    feedback_id = str(uuid.uuid4())
    user_id = current["id"] if current else None
    user_email = current["email"] if current else payload.email

    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO feedback (id, user_id, type, content, email, created_at)
                VALUES (%s, %s, %s, %s, %s, %s)
                """,
                (feedback_id, user_id, payload.type, payload.content, user_email, datetime.now(timezone.utc)),
            )

    # Email sending logic
    try:
        smtp_server = os.environ.get("SMTP_SERVER")
        smtp_port = int(os.environ.get("SMTP_PORT", "587"))
        smtp_user = os.environ.get("SMTP_USER")
        smtp_pass = os.environ.get("SMTP_PASS")
        
        if smtp_server and smtp_user and smtp_pass:
            msg = MIMEMultipart()
            msg['From'] = smtp_user
            msg['To'] = "contact@pixelbond.in"
            msg['Subject'] = f"New {payload.type.replace('_', ' ').title()} from Roastmaster.ai"
            
            body = f"Type: {payload.type}\nContent: {payload.content}\nUser Email: {user_email or 'Not provided'}\nUser ID: {user_id or 'Guest'}"
            msg.attach(MIMEText(body, 'plain'))
            
            with smtplib.SMTP(smtp_server, smtp_port) as server:
                server.starttls()
                server.login(smtp_user, smtp_pass)
                server.send_message(msg)
    except Exception as e:
        logger.error(f"Failed to send feedback email: {e}")

    return {"status": "success", "message": "Feedback submitted successfully"}


@api_router.get("/")
async def root():
    return {"message": "Roastmaster API v1"}


app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get("CORS_ORIGINS", "*").split(","),
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/{full_path:path}", include_in_schema=False)
async def serve_frontend(full_path: str):
    if full_path.startswith("api/"):
        raise HTTPException(status_code=404, detail="Not found")
    index_file = FRONTEND_BUILD_DIR / "index.html"
    candidate = (FRONTEND_BUILD_DIR / full_path).resolve() if full_path else index_file
    try:
        candidate.relative_to(FRONTEND_BUILD_DIR.resolve())
    except ValueError as exc:
        raise HTTPException(status_code=404, detail="Not found") from exc

    if full_path and candidate.exists() and candidate.is_file():
        return FileResponse(candidate)
    if index_file.exists():
        return FileResponse(index_file)
    raise HTTPException(status_code=404, detail="Frontend build not found")
